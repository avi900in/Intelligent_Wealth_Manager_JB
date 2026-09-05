import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Julius Baer Luxury Private Banking Theme
    NAVY = RGBColor(12, 26, 48)          # #0C1A30 (Deep Midnight Navy)
    GOLD = RGBColor(197, 160, 89)        # #C5A059 (JB Warm Champagne Gold)
    DARK_BLUE = RGBColor(20, 42, 74)     # #142A4A (Slate Card Blue)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(244, 246, 249) # #F4F6F9
    BORDER_GRAY = RGBColor(226, 232, 240)# #E2E8F0
    TEXT_MUTED = RGBColor(148, 163, 184) # #94A3B8
    TEXT_DARK = RGBColor(30, 41, 59)     # #1E293B
    ACCENT_RED = RGBColor(225, 29, 72)   # #E11D48 (Risk / Alert)
    ACCENT_GREEN = RGBColor(16, 185, 129)# #10B981 (Success / Compliance)
    ACCENT_PURPLE = RGBColor(139, 92, 246) # #8B5CF6 (AI Innovation)

    blank_layout = prs.slide_layouts[6] # completely blank

    def add_header(slide, title_text, category_text="JULIUS BAER • WEALTH INTELLIGENCE"):
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.15))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = GOLD
        p_cat.space_after = Pt(2)

        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(21)
        p_title.font.bold = True
        p_title.font.color.rgb = NAVY

    def add_footer(slide, current_page, total_pages=12):
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4))
        tf = footer_box.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = f"SingHacks 2026 — JB Pulse: AI-Powered Wealth Intelligence Cockpit | Confidential & Proprietary | Slide {current_page} of {total_pages}"
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 1: Title Slide (Dark Luxury Navy Theme)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.fill.background()

    # Title Card Accent Bar
    accent_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.7), Inches(0.12), Inches(3.9))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = GOLD
    accent_bar.line.fill.background()

    tb1 = s1.shapes.add_textbox(Inches(1.4), Inches(1.7), Inches(10.5), Inches(4.8))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p0 = tf1.paragraphs[0]
    p0.text = "BANK JULIUS BAER • SINGHAX 2026 CHALLENGE"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = GOLD
    p0.space_after = Pt(10)

    p1 = tf1.add_paragraph()
    p1.text = "JB Pulse — Wealth Intelligence Cockpit"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_after = Pt(6)

    p_tagline = tf1.add_paragraph()
    p_tagline.text = "From Intelligence to Wealth"
    p_tagline.font.size = Pt(22)
    p_tagline.font.bold = True
    p_tagline.font.color.rgb = GOLD
    p_tagline.space_after = Pt(14)

    p2 = tf1.add_paragraph()
    p2.text = "Explainable Multi-Agent Advisory with Deterministic Financial Guardrails & Four-Eyes Governance"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(203, 213, 225)
    p2.space_after = Pt(22)

    p3 = tf1.add_paragraph()
    p3.text = "Spec-Driven Engineering • 6 Specialist Agents • Look-Through Risk • 100% LLM Outage Resiliency • Maker-Checker Sign-off"
    p3.font.size = Pt(11.5)
    p3.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 2: The Core Problem & The Private Banking "Advisory Gap"
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "The Private Banking Challenge: The 'Advisory Gap'", "THE PROBLEM STATEMENT")
    add_footer(s2, 2)

    cards_s2 = [
        ("1. Descriptive, Not Advisory", 
         "Current digital tools stop at descriptive reporting ('What is' — valuations, SAA charts, transaction tables).\n\nRMs must manually deduce why portfolios shifted, how macro shocks transmit to multi-sleeve holdings, and what commercial action to take ('What should I do?').",
         ACCENT_RED),
        ("2. Hidden Multi-Sleeve Complexity", 
         "Clients hold segregated accounts (Discretionary, Advisory, Execution, Family Trusts).\n\nCritical risks (13%+ aggregated single-stock concentration, Lombard cross-collateral LTV, PE capital call shortfalls) are invisible in isolated views and only emerge when aggregated across all sleeves.",
         GOLD),
        ("3. RM Cognitive Overload", 
         "Priscilla Ong manages 20 UHNW relationships ($290M+ AUM) across 1,015 holdings, 5 macro shocks, and 28 informal CRM notes.\n\nAnalyzing portfolios and preparing bespoke client meeting packs takes 4+ hours per client, severely capping RM scale and quality of advice.",
         NAVY)
    ]

    for i, (title, desc, color) in enumerate(cards_s2):
        left = Inches(0.8 + i * 3.95)
        top = Inches(1.8)
        w = Inches(3.75)
        h = Inches(4.7)

        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = BORDER_GRAY

        line = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.1))
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()

        tb = s2.shapes.add_textbox(left + Inches(0.25), top + Inches(0.3), w - Inches(0.5), h - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = NAVY
        pt.space_after = Pt(14)

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(12)
        pd.font.color.rgb = TEXT_DARK
        pd.space_after = Pt(10)

    # =========================================================================
    # SLIDE 3: Approach Taken: Spec-Driven Agentic Engineering
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Approach: Spec-Driven Development (From Intelligence to Wealth)", "ENGINEERING METHODOLOGY")
    add_footer(s3, 3)

    spec_phases = [
        ("Phase 1: Deep Domain & Data Analysis",
         "• Comprehensive inspection of private banking data structures across 12 source files (clients, portfolios, 1,015 holdings, credit facilities, commitments, event log, unstructured RM notes).\n• Identified core traps: multi-portfolio mandate flattening and structured product look-through opacity.\n• Mapped high-touch RM commercial workflows for Asia Desks.",
         NAVY),
        ("Phase 2: Multi-LLM Spec Synthesis",
         "• Requirements, constraints, and architecture were rigorously debated and authored collaboratively using Claude 3.7 Sonnet/Opus, ChatGPT-4o/o1, and Google Gemini 2.0 Pro.\n• Formally published as WEALTH_INTELLIGENCE_AGENT_SPEC.md.\n• Enforced architectural principle: Strict separation of deterministic compute from LLM reasoning.",
         GOLD),
        ("Phase 3: Autonomous Agentic Build",
         "• Spec document fed into autonomous agentic coding workflows to implement modular specialist agents.\n• Engineered 6 specialist advisory agents (Rebalancing, Tax Optimization, Life Event Planning, Liquidity/Credit Risk, Market Impact, RM Qualitative Notes).\n• Built standardized JSON contracts, evidence citations, and temporal horizons.",
         ACCENT_PURPLE)
    ]

    for i, (title, desc, color) in enumerate(spec_phases):
        left = Inches(0.8 + i * 3.95)
        top = Inches(1.8)
        w = Inches(3.75)
        h = Inches(4.7)

        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = BORDER_GRAY

        line = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.1))
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()

        tb = s3.shapes.add_textbox(left + Inches(0.25), top + Inches(0.3), w - Inches(0.5), h - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(15)
        pt.font.bold = True
        pt.font.color.rgb = NAVY
        pt.space_after = Pt(12)

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11.5)
        pd.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 4: Current Setup: Agent Swarm, Deterministic Engine & Master Orchestrator
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "Current Setup: Three-Tier Agent Swarm & Deterministic Orchestration", "CURRENT SYSTEM ARCHITECTURE")
    add_footer(s4, 4)

    arch_layers = [
        ("Layer 1: Deterministic Analytics Engine (Zero Hallucination Core)", 
         "Calculates SAA drift %, structured product look-through decomposition, consolidated Lombard LTV trajectory, PE capital call runway, and tax-loss harvesting lots mathematically using pure functional code before any LLM invocation.", NAVY),
        ("Layer 2: Multi-Specialist Advisory Swarm (6 Specialized Micro-Agents)", 
         "• Rebalancing Agent: Single & cross-portfolio mandate breaches, trade sizing, phased execution horizon (2-4 wks)\n• Tax Optimization Agent: Loss harvesting offsets, domicile tax rules, tax year-end window (pre-31 Dec)\n• Life Event Planning Agent: Ring-fences committed cash buffers for dynamic milestone dates (e.g. property, tuition)\n• Liquidity & Credit Risk Agent: Lombard LTV margin buffer remediation (48-72h) & uncalled PE capital reserves\n• Market Impact Agent: Macro event_log transmission channels & tactical hedging tenors (3-6 mos)\n• RM Notes & Qualitative Agent: Unstructured CRM notes, standing exclusions & permanent holding policies", DARK_BLUE),
        ("Layer 3: Master Orchestrator & Synergistic Synthesis Engine", 
         "• Conflict Resolution: Resolves trade tensions (e.g. mandate trim vs. RM legacy holding override)\n• Comingling Opportunities: Clubs multi-agent actions into unified, cost-efficient client proposals\n• Composite Urgency Scoring (0-100) & Point-in-Time Event Filtering: Highlights developments post last RM meeting\n• Meeting Pack Synthesis: Generates client-ready briefings, formal emails, and supervisory audit dossiers", GOLD)
    ]

    for i, (title, desc, color) in enumerate(arch_layers):
        top = Inches(1.68 + i * 1.62)
        box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = BORDER_GRAY

        bar = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(0.12), Inches(1.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        tb = s4.shapes.add_textbox(Inches(1.1), top + Inches(0.12), Inches(11.2), Inches(1.25))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(14)
        pt.font.bold = True
        pt.font.color.rgb = NAVY
        pt.space_after = Pt(4)

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 5: Technology Stack: Current Prototype vs. Proposed Production Deployment
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "Technology Architecture: Current Prototype vs. Production Deployment", "TECH STACK COMPARISON")
    add_footer(s5, 5)

    # Left Column: Current Hackathon Stack
    left_stack = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(5.65), Inches(4.9))
    left_stack.fill.solid()
    left_stack.fill.fore_color.rgb = LIGHT_GRAY
    left_stack.line.color.rgb = BORDER_GRAY

    line_l = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.7), Inches(5.65), Inches(0.1))
    line_l.fill.solid()
    line_l.fill.fore_color.rgb = NAVY
    line_l.line.fill.background()

    tb_l = s5.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.15), Inches(4.5))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "⚡ Current Hackathon Prototype Stack"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(10)

    current_items = [
        "Frontend & UI: Streamlit (Python 3.13) + Julius Baer Bespoke Luxury CSS Theme & Design System.",
        "Agent Framework: Custom Multi-Agent Swarm with Pydantic State Schema & Tool Execution Contracts.",
        "LLM Engine: Multi-Model Gateway (Claude 3.7 / OpenAI GPT-4o / Google Gemini API wrappers).",
        "Deterministic Analytics: Pure functional Python, NumPy, and Pandas calculation modules.",
        "Semantic Search / Vector Store: In-memory ChromaDB semantic index for RM notes and macro event log.",
        "Storage & Repository: CSV/JSON relational repository with hierarchical entity graph (Client → Portfolios → Holdings).",
        "Governance & Security: Streamlit session state approval cache with instant JSON audit log export."
    ]
    for item in current_items:
        p = tf_l.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(5)

    # Right Column: Proposed Enterprise Production Stack
    right_stack = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.7), Inches(4.9))
    right_stack.fill.solid()
    right_stack.fill.fore_color.rgb = DARK_BLUE
    right_stack.line.fill.background()

    line_r = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.7), Inches(0.1))
    line_r.fill.solid()
    line_r.fill.fore_color.rgb = GOLD
    line_r.line.fill.background()

    tb_r = s5.shapes.add_textbox(Inches(7.05), Inches(1.9), Inches(5.2), Inches(4.5))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "🏛️ Proposed Enterprise Production Deployment"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.space_after = Pt(10)

    prod_items = [
        "Frontend: Next.js 15 / React 19 + Radix UI + Julius Baer EAM/RM Portal Micro-frontend integration.",
        "API & Services: FastAPI / Async Python microservices with Ray / Celery Distributed Task Worker Clusters.",
        "Orchestration Engine: LangGraph Enterprise / Temporal Workflow Orchestrator with stateful checkpoints.",
        "LLM Gateway: Dedicated Private Azure OpenAI / AWS Bedrock VPC (Zero Data Retention) + On-Prem Air-Gapped SLM (Llama 3 70B / Mistral Large).",
        "Deterministic Kernel: High-performance C++ / Rust financial calculation engine (microsecond latency).",
        "Data Warehouse: Snowflake Data Cloud + TimescaleDB (time-series snapshots) + Milvus/pgvector Enterprise Vector DB.",
        "Security & Compliance: Enterprise SAML/OIDC SSO, HashiCorp Vault, and Kafka immutable audit ledger for FINMA/MAS compliance."
    ]
    for item in prod_items:
        p = tf_r.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.space_after = Pt(5)

    # =========================================================================
    # SLIDE 6: Enterprise Guardrails: Determinism, Four-Eyes Governance & Fallbacks
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "Enterprise Guardrails: Deterministic Compute, Four-Eyes Control & Fallbacks", "TRUST & RISK CONTROLS")
    add_footer(s6, 6)

    guardrail_cards = [
        ("a) Deterministic Layer (Zero Hallucination)", 
         "• Strict decoupling of compute from reasoning.\n• All drift %, LTV ratios, look-through weights, and cash runways are computed via pure immutable functions.\n• LLMs NEVER perform arithmetic; they only receive pre-verified metric facts as tool outputs and cite source function names and valuation dates.",
         NAVY),
        ("b) Maker-Checker Supervisory Control", 
         "• Maker (RM Priscilla Ong): Reviews recommendations, customizes talking points, queues commercial actions.\n• Checker (Desk Head Marc Guggenheim DH-SG-001): 4-Point Supervisory Audit (KYC, Suitability, Exclusions, Cross-Border Fit).\n• Digital Endorsement Stamp & SHA-256 audit token required before client export.",
         GOLD),
        ("c) Graceful Deterministic Fallback", 
         "• 100% LLM Outage & Latency Resiliency.\n• If LLM APIs are offline, rate-limited, or latent, system automatically fails over to instant deterministic rule synthesis.\n• Produces complete, compliant, professional client recommendations and talking points deterministically without service disruption.",
         ACCENT_GREEN)
    ]

    for i, (title, desc, color) in enumerate(guardrail_cards):
        left = Inches(0.8 + i * 3.95)
        top = Inches(1.8)
        w = Inches(3.75)
        h = Inches(4.7)

        card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = BORDER_GRAY

        line = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.1))
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()

        tb = s6.shapes.add_textbox(left + Inches(0.25), top + Inches(0.3), w - Inches(0.5), h - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(14.5)
        pt.font.bold = True
        pt.font.color.rgb = NAVY
        pt.space_after = Pt(12)

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11.5)
        pd.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 7: Core Innovation - Look-Through Decomposition & Whole-Client Risk
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Innovation: Unmasking Invisible Cross-Portfolio & Derivative Risks", "CORE CAPABILITIES")
    add_footer(s7, 7)

    # 2 Columns
    left_c7 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(5.65), Inches(4.9))
    left_c7.fill.solid()
    left_c7.fill.fore_color.rgb = LIGHT_GRAY
    left_c7.line.color.rgb = BORDER_GRAY

    tb_l7 = s7.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.05), Inches(4.5))
    tf_l7 = tb_l7.text_frame
    tf_l7.word_wrap = True

    p = tf_l7.paragraphs[0]
    p.text = "🔍 Look-Through Derivative Decomposition"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(12)

    bullets_l7 = [
        "Core banking systems only see the product wrapper ('FCN ref Basket A' or 'Asia Autocallable').",
        "Our engine unpacks 'worst-of' structured notes into constituent single equities (e.g. Helios Cloud, Global Energy Majors).",
        "Allocates fractional economic exposure dynamically to underlying asset classes and sectors.",
        "Prevents unmonitored concentration build-up in volatile single names across advisory and execution sleeves."
    ]
    for b in bullets_l7:
        p = tf_l7.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(7)

    right_c7 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.7), Inches(4.9))
    right_c7.fill.solid()
    right_c7.fill.fore_color.rgb = LIGHT_GRAY
    right_c7.line.color.rgb = BORDER_GRAY

    tb_r7 = s7.shapes.add_textbox(Inches(7.1), Inches(1.9), Inches(5.1), Inches(4.5))
    tf_r7 = tb_r7.text_frame
    tf_r7.word_wrap = True

    p = tf_r7.paragraphs[0]
    p.text = "🌐 Whole-Client Risk Aggregation"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(12)

    bullets_r7 = [
        "Aggregate Concentration: 6% stock in Sleeve A + 7% in Sleeve B = 13% client-level breach (flagged as CONC_AGG).",
        "Consolidated Lombard LTV: Combines collateral across all pledged sleeves against credit facility FAC-0001.",
        "Liquidity Pool vs PE Deficits: Pools cash across sleeves to test against uncalled capital calls ($2.5M) and planned outflows.",
        "Macro Shock Transmission: Aggregates client-wide exposure to oil, rates, semiconductors, and currency shifts."
    ]
    for b in bullets_r7:
        p = tf_r7.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(7)

    # =========================================================================
    # SLIDE 8: Deep Dive Case Studies - Cheung Kwok Wing & Hartono Kusuma
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "Client Deep Dive: Cheung Kwok Wing (CL-0012) & Hartono Kusuma (CL-0001)", "WORKED CLIENT CASE STUDY")
    add_footer(s8, 8)

    c1 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(3.7), Inches(4.9))
    c1.fill.solid()
    c1.fill.fore_color.rgb = LIGHT_GRAY
    c1.line.color.rgb = BORDER_GRAY
    tb1 = s8.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(3.3), Inches(4.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "👤 Cheung Kwok Wing (CL-0012)"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = NAVY
    p.space_after = Pt(8)
    for b in [
        "Age 71, retired in Hong Kong ($28.0M AUM).",
        "Draws $1.1M/yr for living expenses.",
        "RM Note: Explicitly stated he will NOT sell bonds at a loss.",
        "Fixed income down -$5.6M due to rate hikes.",
        "Longest bond matures in 2045: waiting for recovery is NOT a plan he can outlive."
    ]:
        p = tf1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(4)

    c2 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.7), Inches(3.7), Inches(4.9))
    c2.fill.solid()
    c2.fill.fore_color.rgb = LIGHT_GRAY
    c2.line.color.rgb = BORDER_GRAY
    tb2 = s8.shapes.add_textbox(Inches(5.0), Inches(1.85), Inches(3.3), Inches(4.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "⚠️ Hartono Kusuma (CL-0001)"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_RED
    p.space_after = Pt(8)
    for b in [
        "Age 34, Indonesian 2nd Gen ($46.6M across 2 portfolios).",
        "41.4% concentrated in Bara Nusantara Energy across sleeves.",
        "Lombard facility LTV reached 68.4% (near 70% margin call threshold).",
        "Cross-portfolio look-through unmasks hidden aggregate risk invisible within individual accounts."
    ]:
        p = tf2.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(4)

    c3 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.7), Inches(3.7), Inches(4.9))
    c3.fill.solid()
    c3.fill.fore_color.rgb = NAVY
    c3.line.fill.background()
    tb3 = s8.shapes.add_textbox(Inches(9.0), Inches(1.85), Inches(3.3), Inches(4.5))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "🎯 Client-Ready Solutions"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = GOLD
    p.space_after = Pt(8)
    for b in [
        "For Cheung Kwok Wing: Barbell reinvestment into 2Y Treasury ladder (5.1% yield) funding $1.1M income without 2045 duration risk.",
        "For Hartono Kusuma: Trim $2.4M of energy holding over 2-4 weeks to pay down Lombard debt (lowering LTV to safe 58%).",
        "Conversation Scripts: Empathetic, psychologically grounded talking points for Priscilla."
    ]:
        p = tf3.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.space_after = Pt(4)

    # =========================================================================
    # SLIDE 9: Deep Dive Case Study - Kim Do-Yoon (CL-0015 — PE Runway)
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "Client Deep Dive: Kim Do-Yoon (CL-0015 — Tech Founder, $15.0M AUM)", "WORKED CLIENT CASE STUDY")
    add_footer(s9, 9)

    cd1 = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(3.7), Inches(4.9))
    cd1.fill.solid()
    cd1.fill.fore_color.rgb = LIGHT_GRAY
    cd1.line.color.rgb = BORDER_GRAY
    tbd1 = s9.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(3.3), Inches(4.5))
    tfd1 = tbd1.text_frame
    tfd1.word_wrap = True
    p = tfd1.paragraphs[0]
    p.text = "💼 The Client Situation"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = NAVY
    p.space_after = Pt(8)
    for b in [
        "Age 38, Tech Entrepreneur ($15.0M AUM).",
        "Wealth created via AI software startup Series C round.",
        "High growth appetite, aggressive risk tolerance score.",
        "Committed $5.0M to Private Equity Growth Fund ($2.5M uncalled capital)."
    ]:
        p = tfd1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(4)

    cd2 = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.7), Inches(3.7), Inches(4.9))
    cd2.fill.solid()
    cd2.fill.fore_color.rgb = LIGHT_GRAY
    cd2.line.color.rgb = BORDER_GRAY
    tbd2 = s9.shapes.add_textbox(Inches(5.0), Inches(1.85), Inches(3.3), Inches(4.5))
    tfd2 = tbd2.text_frame
    tfd2.word_wrap = True
    p = tfd2.paragraphs[0]
    p.text = "🚨 Liquidity Shortfall Risk"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_RED
    p.space_after = Pt(8)
    for b in [
        "Upcoming PE Capital Call: $1.2M due in Q4 2026.",
        "Cash holdings in portfolio: only $340,000 (Deficit of -$860,000).",
        "Portfolio heavily concentrated in illiquid unlisted shares (Aranya Tech) and structured autocallable notes.",
        "Forced fire-sale risk if PE drawdowns accelerate."
    ]:
        p = tfd2.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(4)

    cd3 = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.7), Inches(3.7), Inches(4.9))
    cd3.fill.solid()
    cd3.fill.fore_color.rgb = NAVY
    cd3.line.fill.background()
    tbd3 = s9.shapes.add_textbox(Inches(9.0), Inches(1.85), Inches(3.3), Inches(4.5))
    tfd3 = tbd3.text_frame
    tfd3.word_wrap = True
    p = tfd3.paragraphs[0]
    p.text = "💡 Tax-Aware Rebalancing"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = GOLD
    p.space_after = Pt(8)
    for b in [
        "Tax-Aware Domicile Structuring: Rebalance with zero personal capital gains tax drag.",
        "Structured Liquidity Buffer: Liquidate $1.0M of maturing Asian Bank Autocallables into money market funds.",
        "Collateral Optimization: Pre-approve Lombard line against unlisted shares to establish a $1.5M credit liquidity backstop."
    ]:
        p = tfd3.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.space_after = Pt(4)

    # =========================================================================
    # SLIDE 10: RM Workbench - Whole-Book Prioritization Engine
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "RM Workbench: Whole-Book Prioritization Engine", "WORKBENCH INNOVATION")
    add_footer(s10, 10)

    left_w = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(4.2), Inches(4.9))
    left_w.fill.solid()
    left_w.fill.fore_color.rgb = LIGHT_GRAY
    left_w.line.color.rgb = BORDER_GRAY

    tb_lw = s10.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(3.8), Inches(4.5))
    tf_lw = tb_lw.text_frame
    tf_lw.word_wrap = True
    p = tf_lw.paragraphs[0]
    p.text = "📊 Prioritization Algorithm"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = NAVY
    p.space_after = Pt(12)

    formula_items = [
        "Composite Risk Score (0-100): Combines mandate breach severity, LTV margin call buffer, PE liquidity deficits, and days since last interaction.",
        "Point-in-Time Last Meeting Filter: Focuses strictly on market events post-dating the last RM conversation.",
        "Defensible Ranking: Priscilla knows exactly why Client #1 ranks ahead of Client #20 on her morning dashboard."
    ]
    for b in formula_items:
        p = tf_lw.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    right_w = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(1.7), Inches(7.2), Inches(4.9))
    right_w.fill.solid()
    right_w.fill.fore_color.rgb = LIGHT_GRAY
    right_w.line.color.rgb = BORDER_GRAY

    tb_rw = s10.shapes.add_textbox(Inches(5.5), Inches(1.9), Inches(6.8), Inches(4.5))
    tf_rw = tb_rw.text_frame
    tf_rw.word_wrap = True
    p = tf_rw.paragraphs[0]
    p.text = "🏆 Priscilla's Morning Call Queue (Aug 26, 2026)"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = NAVY
    p.space_after = Pt(10)

    queue_rows = [
        ("Rank 1: Hartono Wijaya Kusuma ($46.6M | CL-0001)", "Score 89 | Lombard LTV 68.4% (Near Margin Call) + 41.4% Energy Conc."),
        ("Rank 2: Cheung Kwok Wing ($28.0M | CL-0012)", "Score 86 | $5.6M Fixed Income Duration Deficit vs $1.1M/yr Cash Needs"),
        ("Rank 3: Kim Do-Yoon ($15.0M | CL-0015)", "Score 84 | PE Capital Call Shortfall (-$860k) + Tech Concentration"),
        ("Rank 4: Fong Enterprises FO ($87.9M | CL-0017)", "Score 78 | Cross-Portfolio Energy Overweight + SAA Drift Breaches"),
        ("Rank 5-20: Monitored Stable Book", "Automated tracking; alerts triggered only on threshold breaches")
    ]
    for client, reason in queue_rows:
        p = tf_rw.add_paragraph()
        p.text = f"▶ {client}\n   {reason}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(5)

    # =========================================================================
    # SLIDE 11: End-to-End Workflow: From Signal to Client Meeting Pack
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, "End-to-End Workflow: From Signal to Supervisory Endorsement", "USER EXPERIENCE")
    add_footer(s11, 11)

    steps = [
        ("Step 1: Morning Priority Queue", "Priscilla logs in → Book Overview → Top 5 prioritized clients flagged by composite urgency scoring."),
        ("Step 2: Client 360 & Stress Testing", "Unpacks holdings, 5-snapshot trends, cross-portfolio LTV & PE liquidity runways in interactive lab."),
        ("Step 3: Specialist Agent Action Deck", "Reviews multi-agent recommendations, customizes talking points, approves rebalancing orders."),
        ("Step 4: Supervisory 4-Point Audit", "Desk Head Marc Guggenheim executes four-eyes suitability audit and applies digital approval stamp."),
        ("Step 5: Client Meeting Pack Export", "Exports branded Briefing Notes, Client Discussion Agenda, or Formal Advisory Email in 1 click.")
    ]

    for i, (title, desc) in enumerate(steps):
        top = Inches(1.7 + i * 0.98)
        box = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(0.85))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = BORDER_GRAY

        badge = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(2.2), Inches(0.85))
        badge.fill.solid()
        badge.fill.fore_color.rgb = NAVY if i < 4 else GOLD
        badge.line.fill.background()

        tbb = s11.shapes.add_textbox(Inches(0.9), top + Inches(0.22), Inches(2.0), Inches(0.4))
        tfb = tbb.text_frame
        p = tfb.paragraphs[0]
        p.text = f"STAGE 0{i+1}"
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE

        tb = s11.shapes.add_textbox(Inches(3.2), top + Inches(0.12), Inches(9.1), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(13)
        pt.font.bold = True
        pt.font.color.rgb = NAVY

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 12: Strategic Impact & Why Julius Baer Wins
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    bg12 = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg12.fill.solid()
    bg12.fill.fore_color.rgb = NAVY
    bg12.line.fill.background()

    header_box = s12.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.1))
    tf = header_box.text_frame
    p_cat = tf.paragraphs[0]
    p_cat.text = "STRATEGIC IMPACT • SINGHAX 2026".upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = GOLD
    p_title = tf.add_paragraph()
    p_title.text = "Empowering the RM: The Future of Private Wealth Advisory"
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE

    pillars = [
        ("⏱️ 10x Advisory Productivity", "Cuts RM meeting preparation time from 4 hours to 15 minutes per client. RMs can scale from 20 to 50+ HNW relationships without sacrificing bespoke intimacy."),
        ("🛡️ Zero-Defect Risk Governance", "Deterministic guardrails eliminate compliance drift, unmonitored margin calls, and unauthorized trades. Every action has an immutable audit trail."),
        ("💡 Measurable Alpha & Fee Retention", "Proactive tax-loss harvesting and liquidity management turn defensive portfolio management into high-value advisory dialogue."),
        ("🤝 Trust & Human Centrality", "AI amplifies rather than replaces the Relationship Manager. Julius Baer retains its core heritage: high-touch human relationships powered by world-class intelligence.")
    ]

    for i, (title, desc) in enumerate(pillars):
        col = i % 2
        row = i // 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.9 + row * 2.5)
        w = Inches(5.75)
        h = Inches(2.2)

        card = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BLUE
        card.line.color.rgb = RGBColor(30, 58, 95)

        tb = s12.shapes.add_textbox(left + Inches(0.3), top + Inches(0.2), w - Inches(0.6), h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = GOLD
        pt.space_after = Pt(8)

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(12)
        pd.font.color.rgb = WHITE

    # Save Presentation
    output_path = "/Users/avijit/Documents/Avi/Claude Code AI Hackathon/App Builder/julius_baer_wealth_intelligence_pitch.pptx"
    prs.save(output_path)
    print(f"Presentation successfully created at: {output_path}")

if __name__ == "__main__":
    create_deck()
